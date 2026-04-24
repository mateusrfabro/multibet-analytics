"""
Gera PPTX executivo (1 slide) - Diagnostico Play4Tune.
Audiencia: Head (Castrin), CGO (Conson), CTO (Gabriel).
Linguagem: direta, sem jargao tecnico.
"""
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Cores ----------
COR_HEADER       = RGBColor(0x0B, 0x2B, 0x40)   # azul petroleo
COR_SUB_HEADER   = RGBColor(0x37, 0x47, 0x56)   # cinza chumbo
COR_DESTAQUE     = RGBColor(0xC6, 0x28, 0x28)   # vermelho perda
COR_POSITIVO     = RGBColor(0x1B, 0x5E, 0x20)   # verde
COR_TXT          = RGBColor(0x25, 0x2B, 0x33)   # quase preto
COR_CINZA        = RGBColor(0x70, 0x78, 0x82)   # cinza suporte
COR_BG_CARD      = RGBColor(0xF5, 0xF7, 0xFA)   # bg claro
COR_DIVISORIA    = RGBColor(0xD0, 0xD7, 0xDE)

# ---------- Helpers ----------
def add_text(slide, left, top, width, height, text, *,
             size=11, bold=False, color=COR_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rich_text(slide, left, top, width, height, segments, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """segments = list of (text, {size, bold, color}) per paragraph. Cada segmento vira um paragrafo."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    for i, seg in enumerate(segments):
        text, opts = seg
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(opts.get("space_after", 2))
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(opts.get("size", 11))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = opts.get("color", COR_TXT)
    return tb


def add_rect(slide, left, top, width, height, fill=COR_BG_CARD, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    # desabilita sombra
    return shape


def add_line(slide, x1, y1, x2, y2, color=COR_DIVISORIA, weight=1.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


# ---------- Build ----------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # ---------- Fundo branco garantido ----------
    # (blank layout ja e branco por default)

    # ======================================================
    # CABECALHO (barra superior)
    # ======================================================
    header_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0), Inches(0),
                                         prs.slide_width, Inches(0.9))
    header_rect.fill.solid()
    header_rect.fill.fore_color.rgb = COR_HEADER
    header_rect.line.fill.background()

    add_text(slide, Inches(0.35), Inches(0.12), Inches(10), Inches(0.45),
             "PLAY4TUNE — DIAGNOSTICO DE FUNIL DE AQUISICAO",
             size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(0.35), Inches(0.52), Inches(10), Inches(0.3),
             "Causa raiz: instabilidade do gateway de pagamento esta segurando o crescimento",
             size=12, color=RGBColor(0xCD, 0xD9, 0xE3))

    # Data/periodo no topo direito
    add_text(slide, Inches(10.2), Inches(0.2), Inches(3), Inches(0.3),
             "Periodo: 01 a 22/Abr/2026",
             size=10, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(10.2), Inches(0.48), Inches(3), Inches(0.3),
             "Executado em 23/04/2026",
             size=9, color=RGBColor(0xCD, 0xD9, 0xE3), align=PP_ALIGN.RIGHT)

    # ======================================================
    # 3 COLUNAS
    # ======================================================
    # Coluna geometry:
    col_top = Inches(1.15)
    col_h   = Inches(5.4)
    col_w   = Inches(4.25)
    col_gap = Inches(0.08)
    col1_x  = Inches(0.22)
    col2_x  = Inches(0.22 + 4.25 + 0.08)
    col3_x  = Inches(0.22 + (4.25 + 0.08) * 2)

    # Card backgrounds
    add_rect(slide, col1_x, col_top, col_w, col_h, fill=COR_BG_CARD)
    add_rect(slide, col2_x, col_top, col_w, col_h, fill=COR_BG_CARD)
    add_rect(slide, col3_x, col_top, col_w, col_h, fill=COR_BG_CARD)

    # ---------- COLUNA 1: O QUADRO ----------
    y = col_top + Inches(0.1)
    add_text(slide, col1_x + Inches(0.2), y, col_w - Inches(0.4), Inches(0.35),
             "1. O QUADRO — 22 DIAS", size=14, bold=True, color=COR_HEADER)

    # Funil numerico (destaque grande)
    y = col_top + Inches(0.55)
    funnel_data = [
        ("2.393",   "cadastros novos",            COR_TXT),
        ("1.677",   "tentaram depositar",         COR_TXT),
        ("648",     "depositos aprovados",        COR_SUB_HEADER),
        ("512",     "viraram primeiro deposito",  COR_POSITIVO),
    ]
    for numero, label, cor in funnel_data:
        add_text(slide, col1_x + Inches(0.25), y, Inches(1.4), Inches(0.45),
                 numero, size=22, bold=True, color=cor, align=PP_ALIGN.RIGHT)
        add_text(slide, col1_x + Inches(1.75), y + Inches(0.08), Inches(2.4), Inches(0.4),
                 label, size=11, color=COR_TXT)
        y += Inches(0.58)

    # Taxa conversao final destacada
    y += Inches(0.1)
    add_text(slide, col1_x + Inches(0.25), y, col_w - Inches(0.5), Inches(0.3),
             "Conversao cadastro → deposito: 21%",
             size=11, bold=True, color=COR_SUB_HEADER)
    y += Inches(0.35)

    # Linha divisoria
    add_line(slide,
             col1_x + Inches(0.3), y,
             col1_x + col_w - Inches(0.3), y)
    y += Inches(0.15)

    # Contexto tracking
    add_rich_text(slide, col1_x + Inches(0.25), y, col_w - Inches(0.5), Inches(1.6),
                  [
                      ("Origem do trafego", {"size": 11, "bold": True, "color": COR_SUB_HEADER, "space_after": 4}),
                      ("86% vem de Meta (Facebook 64% + Instagram 22%).", {"size": 10.5, "color": COR_TXT, "space_after": 2}),
                      ("Tracking funciona (89% dos cadastros com fbclid).", {"size": 10.5, "color": COR_TXT, "space_after": 8}),
                      ("Melhor dia: 07/Abr", {"size": 11, "bold": True, "color": COR_POSITIVO, "space_after": 2}),
                      ("189 cadastros, 42 primeiros depositos, 0 falhas de gateway.", {"size": 10.5, "color": COR_TXT}),
                  ])

    # ---------- COLUNA 2: ONDE O DINHEIRO ESCAPA ----------
    y = col_top + Inches(0.1)
    add_text(slide, col2_x + Inches(0.2), y, col_w - Inches(0.4), Inches(0.35),
             "2. ONDE O DINHEIRO ESTA ESCAPANDO", size=14, bold=True, color=COR_DESTAQUE)

    # Destaque grande: 864 EXPIRARAM
    y = col_top + Inches(0.6)
    add_text(slide, col2_x + Inches(0.25), y, col_w - Inches(0.5), Inches(0.7),
             "864", size=44, bold=True, color=COR_DESTAQUE)
    add_text(slide, col2_x + Inches(0.25), y + Inches(0.8), col_w - Inches(0.5), Inches(0.35),
             "depositos expiraram no meio do caminho",
             size=12, bold=True, color=COR_SUB_HEADER)
    add_text(slide, col2_x + Inches(0.25), y + Inches(1.1), col_w - Inches(0.5), Inches(0.3),
             "(51% de TODAS as tentativas de deposito)",
             size=10.5, color=COR_CINZA)

    y = col_top + Inches(2.25)
    add_line(slide, col2_x + Inches(0.3), y, col2_x + col_w - Inches(0.3), y)
    y += Inches(0.15)

    # Perda quantificada
    add_rich_text(slide, col2_x + Inches(0.25), y, col_w - Inches(0.5), Inches(2.8),
                  [
                      ("251 jogadores nunca voltaram", {"size": 13, "bold": True, "color": COR_DESTAQUE, "space_after": 2}),
                      ("Registraram, tentaram depositar, foi recusado — e nunca mais apareceram.", {"size": 10.5, "color": COR_TXT, "space_after": 10}),
                      ("R$ 7.800 em intencao de deposito perdida", {"size": 13, "bold": True, "color": COR_DESTAQUE, "space_after": 2}),
                      ("441 mil PKR que jogadores tentaram depositar mas nao conseguiram.", {"size": 10.5, "color": COR_TXT, "space_after": 10}),
                      ("Pior horario: 22h-23h (Paquistao)", {"size": 11, "bold": True, "color": COR_SUB_HEADER, "space_after": 2}),
                      ("Aprovacao cai de 50% (horario calmo) para 25% no pico. O gateway nao aguenta o volume.", {"size": 10.5, "color": COR_TXT, "space_after": 10}),
                      ("Por que isso acontece?", {"size": 11, "bold": True, "color": COR_SUB_HEADER, "space_after": 2}),
                      ("Quem conclui o deposito, conclui em menos de 1 minuto. Quando passa disso, ja era.", {"size": 10.5, "color": COR_TXT}),
                  ])

    # ---------- COLUNA 3: O QUE FAZER ----------
    y = col_top + Inches(0.1)
    add_text(slide, col3_x + Inches(0.2), y, col_w - Inches(0.4), Inches(0.35),
             "3. O QUE FAZER — PRIORIDADE", size=14, bold=True, color=COR_POSITIVO)

    y = col_top + Inches(0.55)

    acoes = [
        ("1", "Ativar gateway reserva ja",
         "Um 2o gateway (OKExPay) foi cadastrado em 23/04, mas nao esta operando. Colocar no ar com roteamento inteligente para 20h-23h (Paquistao).",
         "CTO / Engenharia  •  esta semana"),
        ("2", "Deslocar campanhas Meta",
         "Disparar ads entre 14h-20h PKT, quando o gateway funciona bem. Ganho estimado: +10 a +15 pontos percentuais de aprovacao, sem mexer em backend.",
         "CGO / Trafego  •  proxima semana"),
        ("3", "Recuperar deposito em 90s",
         "SMS / WhatsApp automatico 60-90 segundos apos tentativa nao concluida: \"seu deposito esta pendente, clique aqui\". Recuperacao estimada: +50 FTDs/mes.",
         "Produto / CRM  •  2 semanas"),
    ]

    for num, title, desc, owner in acoes:
        # Numerador
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        col3_x + Inches(0.2), y + Inches(0.03),
                                        Inches(0.38), Inches(0.38))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COR_POSITIVO
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = "Calibri"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Titulo
        add_text(slide, col3_x + Inches(0.7), y, col_w - Inches(0.9), Inches(0.35),
                 title, size=12, bold=True, color=COR_HEADER)

        # Descricao
        add_text(slide, col3_x + Inches(0.7), y + Inches(0.32), col_w - Inches(0.9), Inches(0.8),
                 desc, size=10, color=COR_TXT)

        # Dono/prazo
        add_text(slide, col3_x + Inches(0.7), y + Inches(1.1), col_w - Inches(0.9), Inches(0.3),
                 owner, size=9.5, bold=True, color=COR_CINZA)

        y += Inches(1.55)

    # ======================================================
    # RODAPE - O QUE FALTA VALIDAR
    # ======================================================
    rodape_y = Inches(6.65)
    rodape_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0), rodape_y,
                                         prs.slide_width, Inches(0.85))
    rodape_rect.fill.solid()
    rodape_rect.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF6)
    rodape_rect.line.fill.background()

    add_text(slide, Inches(0.35), rodape_y + Inches(0.08), Inches(7), Inches(0.3),
             "O que ainda falta validar",
             size=10.5, bold=True, color=COR_HEADER)
    add_text(slide, Inches(0.35), rodape_y + Inches(0.33), Inches(9), Inches(0.25),
             "•  Custo Meta Ads / CPA real por campanha (pull 24/Abr)    •  Diagnostico com o provedor do gateway (BCTYSO) sobre saturacao 20-23h PKT",
             size=9.5, color=COR_SUB_HEADER)
    add_text(slide, Inches(0.35), rodape_y + Inches(0.55), Inches(9), Inches(0.25),
             "•  Confirmar se o gateway reserva (OKExPay, ativo desde 23/04 01:16) esta recebendo trafego",
             size=9.5, color=COR_SUB_HEADER)

    add_text(slide, Inches(9.8), rodape_y + Inches(0.25), Inches(3.4), Inches(0.3),
             "Elaborado: Mateus Fabro",
             size=9, bold=True, color=COR_HEADER, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(9.8), rodape_y + Inches(0.48), Inches(3.4), Inches(0.25),
             "Fonte: banco Play4Tune  •  Filtros: usuarios de teste + farming excluidos",
             size=8.5, color=COR_CINZA, align=PP_ALIGN.RIGHT)

    # ======================================================
    # SALVAR
    # ======================================================
    out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "reports")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "diagnostico_play4tune_executivo_2026-04-23.pptx")
    prs.save(out_path)
    print(f"OK. PPTX gerado em: {out_path}")
    return out_path


if __name__ == "__main__":
    build()
